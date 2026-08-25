"""Document extraction.

Everything here answers one question: *where exactly did this text come from?* Page
numbers, sheet names, slide numbers and per-word bounding boxes are captured alongside the
text, because the product's whole evidence guarantee rests on being able to open a source
at the exact passage that was cited.
"""

from __future__ import annotations

import io
import re
import zipfile
from dataclasses import dataclass, field
from typing import Any

import pymupdf

from .config import settings

# Coordinates are stored normalised to 0..1 of page width/height so a highlight survives
# whatever zoom level or render width the viewer happens to use.
Box = dict[str, Any]


@dataclass
class Page:
    page_number: int
    text: str
    width: float | None = None
    height: float | None = None
    sheet_name: str | None = None
    slide_number: int | None = None
    ocr_applied: bool = False
    ocr_confidence: float | None = None
    word_boxes: list[Box] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        return {
            "pageNumber": self.page_number,
            "text": self.text,
            "width": self.width,
            "height": self.height,
            "sheetName": self.sheet_name,
            "slideNumber": self.slide_number,
            "ocrApplied": self.ocr_applied,
            "ocrConfidence": self.ocr_confidence,
            "wordBoxes": self.word_boxes,
        }


@dataclass
class Extraction:
    document_type: str
    pages: list[Page]
    title: str | None = None
    metadata: dict[str, Any] = field(default_factory=dict)
    is_scanned: bool = False
    is_signed: bool = False
    is_encrypted: bool = False
    has_macros: bool = False
    has_extractable_text: bool = True
    ocr_applied: bool = False
    ocr_confidence: float | None = None
    media_count: int = 0
    page_sizes: list[dict[str, float]] = field(default_factory=list)
    removed_active_content: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        return {
            "documentType": self.document_type,
            "pages": [p.to_dict() for p in self.pages],
            "pageCount": len(self.pages),
            "title": self.title,
            "metadata": self.metadata,
            "isScanned": self.is_scanned,
            "isSigned": self.is_signed,
            "isEncrypted": self.is_encrypted,
            "hasMacros": self.has_macros,
            "hasExtractableText": self.has_extractable_text,
            "ocrApplied": self.ocr_applied,
            "ocrConfidence": self.ocr_confidence,
            "mediaCount": self.media_count,
            "pageSizes": self.page_sizes,
            "removedActiveContent": self.removed_active_content,
            "warnings": self.warnings,
        }


# ---------------------------------------------------------------------------
# Active content neutralisation
# ---------------------------------------------------------------------------

_SCRIPT_RE = re.compile(r"<script\b[^>]*>.*?</script\s*>", re.IGNORECASE | re.DOTALL)
_STYLE_RE = re.compile(r"<style\b[^>]*>.*?</style\s*>", re.IGNORECASE | re.DOTALL)
_IFRAME_RE = re.compile(r"<iframe\b[^>]*>.*?</iframe\s*>", re.IGNORECASE | re.DOTALL)
_EVENT_RE = re.compile(r"\son[a-z]+\s*=\s*(\"[^\"]*\"|'[^']*')", re.IGNORECASE)
_URI_RE = re.compile(r"\b(?:javascript|vbscript|data):[^\s\"'<>]+", re.IGNORECASE)
_TAG_RE = re.compile(r"<[^>]+>")
_FENCE_RE = re.compile(r"<\|(?:im_start|im_end|system|endoftext)\|>|\[/?(?:INST|SYSTEM|SYS)\]", re.IGNORECASE)


def strip_active_content(text: str) -> tuple[str, list[str]]:
    """Removes executable and instruction-fence content, keeping the readable text.

    The original bytes are never modified; this only affects what gets indexed and quoted,
    so a document that hides a script or a fake system prompt cannot influence anything
    downstream while its legitimate prose stays intact.
    """
    removed: list[str] = []
    out = text

    for name, pattern in (
        ("script", _SCRIPT_RE),
        ("style", _STYLE_RE),
        ("iframe", _IFRAME_RE),
        ("event-handler", _EVENT_RE),
        ("active-uri", _URI_RE),
        ("prompt-fence", _FENCE_RE),
    ):
        new_out, count = pattern.subn(" ", out)
        if count:
            removed.append(f"{name} x{count}")
            out = new_out

    return out, removed


def _normalise(text: str) -> str:
    """Collapses runs of blank lines without destroying line structure.

    Line breaks are load-bearing: heading detection reads them, so they cannot be flattened.
    """
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------


def extract_pdf(data: bytes, max_pages: int, force_ocr: bool = False, password: str | None = None) -> Extraction:
    warnings: list[str] = []
    is_encrypted = False

    try:
        doc = pymupdf.open(stream=data, filetype="pdf")
    except Exception as exc:  # noqa: BLE001 - surfaced to the user verbatim
        raise ValueError(f"This PDF could not be opened: {exc}") from exc

    if doc.needs_pass:
        is_encrypted = True
        if not password or not doc.authenticate(password):
            doc.close()
            raise ValueError(
                "This PDF is password protected. Supply the password to extract its contents."
            )

    is_signed = _pdf_is_signed(doc)
    media_count = 0
    pages: list[Page] = []
    page_sizes: list[dict[str, float]] = []
    ocr_pages = 0
    ocr_confidences: list[float] = []

    total = min(doc.page_count, max_pages)
    if doc.page_count > max_pages:
        warnings.append(
            f"Only the first {max_pages} of {doc.page_count} pages were indexed (page limit)."
        )

    for index in range(total):
        page = doc[index]
        rect = page.rect
        page_sizes.append({"w": float(rect.width), "h": float(rect.height)})
        media_count += len(page.get_images(full=True))

        raw_text = page.get_text("text") or ""
        needs_ocr = force_ocr or len(raw_text.strip()) < settings.ocr_text_threshold

        word_boxes: list[Box] = []
        ocr_applied = False
        ocr_confidence: float | None = None

        if needs_ocr and settings.ocr_enabled and ocr_pages < settings.ocr_max_pages:
            ocr_text, ocr_boxes, ocr_confidence = _ocr_page(page)
            if ocr_text.strip():
                raw_text = ocr_text
                word_boxes = ocr_boxes
                ocr_applied = True
                ocr_pages += 1
                if ocr_confidence is not None:
                    ocr_confidences.append(ocr_confidence)

        if not ocr_applied:
            word_boxes = _pdf_word_boxes(page, rect)

        cleaned, removed = strip_active_content(raw_text)
        if removed:
            warnings.append(f"Page {index + 1}: neutralised {', '.join(removed)}.")

        pages.append(
            Page(
                page_number=index + 1,
                text=_normalise(cleaned),
                width=float(rect.width),
                height=float(rect.height),
                ocr_applied=ocr_applied,
                ocr_confidence=ocr_confidence,
                word_boxes=word_boxes,
            )
        )

    meta = dict(doc.metadata or {})
    title = (meta.get("title") or "").strip() or None
    doc.close()

    text_chars = sum(len(p.text) for p in pages)
    # "Scanned" means the document as a whole carried no usable text layer, which is what
    # decides whether a corrected edition has to be rebuilt from OCR.
    is_scanned = ocr_pages > 0 and ocr_pages >= max(1, len(pages) // 2)

    return Extraction(
        document_type="pdf",
        pages=pages,
        title=title,
        metadata={k: v for k, v in meta.items() if isinstance(v, str)},
        is_scanned=is_scanned,
        is_signed=is_signed,
        is_encrypted=is_encrypted,
        has_extractable_text=text_chars > 0,
        ocr_applied=ocr_pages > 0,
        ocr_confidence=(sum(ocr_confidences) / len(ocr_confidences)) if ocr_confidences else None,
        media_count=media_count,
        page_sizes=page_sizes,
        warnings=warnings,
    )


def _pdf_is_signed(doc: pymupdf.Document) -> bool:
    """Detects a signature field.

    A false positive here is harmless (the user sees a conservative notice); a false
    negative would let the product imply a signature survived an edit, which it must never do.
    """
    try:
        for page in doc:
            for widget in page.widgets() or []:
                if widget.field_type_string and "signature" in widget.field_type_string.lower():
                    return True
    except Exception:  # noqa: BLE001 - detection must never fail the extraction
        pass
    try:
        # Fall back to the raw catalogue: an incremental-update signature may not surface
        # as a widget.
        xref_count = doc.xref_length()
        for xref in range(1, min(xref_count, 3000)):
            obj = doc.xref_object(xref, compressed=False) or ""
            if "/Type/Sig" in obj or "/Type /Sig" in obj or "/ByteRange" in obj:
                return True
    except Exception:  # noqa: BLE001
        pass
    return False


def _pdf_word_boxes(page: pymupdf.Page, rect: pymupdf.Rect) -> list[Box]:
    """Per-word rectangles, normalised to 0..1, used to highlight a cited passage."""
    boxes: list[Box] = []
    width = float(rect.width) or 1.0
    height = float(rect.height) or 1.0
    try:
        for x0, y0, x1, y1, word, *_ in page.get_text("words"):
            if not word.strip():
                continue
            boxes.append(
                {
                    "t": word,
                    "x": round(x0 / width, 5),
                    "y": round(y0 / height, 5),
                    "w": round((x1 - x0) / width, 5),
                    "h": round((y1 - y0) / height, 5),
                }
            )
    except Exception:  # noqa: BLE001 - a page without a text layer simply has no boxes
        return []
    return boxes


def _ocr_page(page: pymupdf.Page) -> tuple[str, list[Box], float | None]:
    """OCRs a rendered page and returns text, word boxes and a mean confidence."""
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        return "", [], None

    try:
        zoom = settings.ocr_dpi / 72.0
        pixmap = page.get_pixmap(matrix=pymupdf.Matrix(zoom, zoom), alpha=False)
        image = Image.open(io.BytesIO(pixmap.tobytes("png")))

        data = pytesseract.image_to_data(
            image, lang=settings.ocr_language, output_type=pytesseract.Output.DICT
        )
    except Exception:  # noqa: BLE001 - OCR is best-effort; failure falls back to no text
        return "", [], None

    width = float(pixmap.width) or 1.0
    height = float(pixmap.height) or 1.0

    boxes: list[Box] = []
    confidences: list[float] = []
    lines: dict[tuple[int, int, int], list[str]] = {}

    for i, word in enumerate(data.get("text", [])):
        if not word or not word.strip():
            continue
        try:
            confidence = float(data["conf"][i])
        except (KeyError, ValueError, IndexError):
            confidence = -1.0
        # Tesseract reports -1 for non-text regions; those carry no information.
        if confidence >= 0:
            confidences.append(confidence)

        key = (data["block_num"][i], data["par_num"][i], data["line_num"][i])
        lines.setdefault(key, []).append(word)

        boxes.append(
            {
                "t": word,
                "x": round(data["left"][i] / width, 5),
                "y": round(data["top"][i] / height, 5),
                "w": round(data["width"][i] / width, 5),
                "h": round(data["height"][i] / height, 5),
            }
        )

    text = "\n".join(" ".join(words) for _, words in sorted(lines.items()))
    mean_confidence = (sum(confidences) / len(confidences) / 100.0) if confidences else None
    return text, boxes, mean_confidence


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------


def extract_docx(data: bytes) -> Extraction:
    import docx

    warnings: list[str] = []
    try:
        document = docx.Document(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        raise ValueError(f"This DOCX could not be opened: {exc}") from exc

    has_macros = _zip_contains(data, ("word/vbaProject.bin",))
    if has_macros:
        warnings.append("This document contains macros. They are never executed or carried forward.")

    blocks: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = (paragraph.style.name or "").lower() if paragraph.style else ""
        # Heading styles are re-emitted as their own line so structure detection sees them
        # even when the author used styles rather than numbering.
        blocks.append(text)
        if "heading" in style and not text.endswith((".", ":", ";")):
            blocks.append("")

    for table in document.tables:
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            blocks.append("")
            blocks.extend(rows)
            blocks.append("")

    text = "\n".join(blocks)
    cleaned, removed = strip_active_content(text)

    core = document.core_properties
    metadata = {
        "author": core.author or "",
        "created": core.created.isoformat() if core.created else "",
        "modified": core.modified.isoformat() if core.modified else "",
        "revision": str(core.revision or ""),
    }

    media_count = _zip_media_count(data, "word/media/")
    # DOCX has no intrinsic pagination; the whole body is one logical page and the locator
    # falls back to heading path plus paragraph index.
    return Extraction(
        document_type="docx",
        pages=[Page(page_number=1, text=_normalise(cleaned))],
        title=(core.title or "").strip() or None,
        metadata=metadata,
        has_macros=has_macros,
        has_extractable_text=bool(cleaned.strip()),
        media_count=media_count,
        page_sizes=[],
        removed_active_content=removed,
        warnings=warnings,
    )


# ---------------------------------------------------------------------------
# XLSX / CSV
# ---------------------------------------------------------------------------


def extract_xlsx(data: bytes) -> Extraction:
    import openpyxl

    warnings: list[str] = []
    has_macros = _zip_contains(data, ("xl/vbaProject.bin",))
    if has_macros:
        warnings.append("This workbook contains macros. They are never executed or carried forward.")

    try:
        # data_only=False keeps formulas visible, which is what a reviewer needs to see;
        # a computed value alone hides how a figure was derived.
        workbook = openpyxl.load_workbook(io.BytesIO(data), data_only=False, read_only=True)
    except Exception as exc:  # noqa: BLE001
        raise ValueError(f"This workbook could not be opened: {exc}") from exc

    pages: list[Page] = []
    for index, sheet_name in enumerate(workbook.sheetnames):
        sheet = workbook[sheet_name]
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            values = ["" if v is None else str(v) for v in row]
            if any(v.strip() for v in values):
                rows.append(" | ".join(values))
        if not rows:
            continue
        cleaned, _ = strip_active_content("\n".join(rows))
        pages.append(Page(page_number=index + 1, text=_normalise(cleaned), sheet_name=sheet_name))

    workbook.close()

    return Extraction(
        document_type="xlsx",
        pages=pages,
        metadata={"sheets": len(pages)},
        has_macros=has_macros,
        has_extractable_text=any(p.text.strip() for p in pages),
        warnings=warnings,
    )


def extract_csv(data: bytes) -> Extraction:
    text = _decode_text(data)
    cleaned, removed = strip_active_content(text)
    return Extraction(
        document_type="csv",
        pages=[Page(page_number=1, text=_normalise(cleaned), sheet_name="Sheet1")],
        has_extractable_text=bool(cleaned.strip()),
        removed_active_content=removed,
    )


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def extract_pptx(data: bytes) -> Extraction:
    from pptx import Presentation

    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        raise ValueError(f"This presentation could not be opened: {exc}") from exc

    pages: list[Page] = []
    media_count = _zip_media_count(data, "ppt/media/")

    for index, slide in enumerate(presentation.slides):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs).strip()
                    if line:
                        parts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        # Speaker notes are part of the record and often carry the substantive detail.
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"Notes: {notes}")

        cleaned, _ = strip_active_content("\n".join(parts))
        pages.append(
            Page(
                page_number=index + 1,
                text=_normalise(cleaned),
                slide_number=index + 1,
                width=float(presentation.slide_width or 0) / 12700.0,
                height=float(presentation.slide_height or 0) / 12700.0,
            )
        )

    return Extraction(
        document_type="pptx",
        pages=pages,
        metadata={"slides": len(pages)},
        has_extractable_text=any(p.text.strip() for p in pages),
        media_count=media_count,
    )


# ---------------------------------------------------------------------------
# Text / Markdown / HTML / images
# ---------------------------------------------------------------------------


def extract_text(data: bytes, document_type: str) -> Extraction:
    text = _decode_text(data)
    cleaned, removed = strip_active_content(text)
    if document_type == "html":
        cleaned = _TAG_RE.sub(" ", cleaned)
        cleaned = (
            cleaned.replace("&nbsp;", " ")
            .replace("&amp;", "&")
            .replace("&lt;", "<")
            .replace("&gt;", ">")
            .replace("&quot;", '"')
        )
    return Extraction(
        document_type=document_type,
        pages=[Page(page_number=1, text=_normalise(cleaned))],
        has_extractable_text=bool(cleaned.strip()),
        removed_active_content=removed,
    )


def extract_image(data: bytes) -> Extraction:
    """Images carry no text layer, so OCR is the only path to content."""
    if not settings.ocr_enabled:
        return Extraction(
            document_type="image",
            pages=[],
            has_extractable_text=False,
            warnings=["OCR is disabled, so no text could be read from this image."],
        )

    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        return Extraction(
            document_type="image",
            pages=[],
            has_extractable_text=False,
            warnings=["OCR is not installed in this worker, so no text could be read."],
        )

    image = Image.open(io.BytesIO(data))
    data_dict = pytesseract.image_to_data(
        image, lang=settings.ocr_language, output_type=pytesseract.Output.DICT
    )

    width = float(image.width) or 1.0
    height = float(image.height) or 1.0
    boxes: list[Box] = []
    confidences: list[float] = []
    lines: dict[tuple[int, int, int], list[str]] = {}

    for i, word in enumerate(data_dict.get("text", [])):
        if not word or not word.strip():
            continue
        try:
            confidence = float(data_dict["conf"][i])
        except (KeyError, ValueError, IndexError):
            confidence = -1.0
        if confidence >= 0:
            confidences.append(confidence)
        key = (data_dict["block_num"][i], data_dict["par_num"][i], data_dict["line_num"][i])
        lines.setdefault(key, []).append(word)
        boxes.append(
            {
                "t": word,
                "x": round(data_dict["left"][i] / width, 5),
                "y": round(data_dict["top"][i] / height, 5),
                "w": round(data_dict["width"][i] / width, 5),
                "h": round(data_dict["height"][i] / height, 5),
            }
        )

    text = "\n".join(" ".join(words) for _, words in sorted(lines.items()))
    confidence = (sum(confidences) / len(confidences) / 100.0) if confidences else None

    return Extraction(
        document_type="image",
        pages=[
            Page(
                page_number=1,
                text=_normalise(text),
                width=float(image.width),
                height=float(image.height),
                ocr_applied=True,
                ocr_confidence=confidence,
                word_boxes=boxes,
            )
        ],
        is_scanned=True,
        has_extractable_text=bool(text.strip()),
        ocr_applied=True,
        ocr_confidence=confidence,
        page_sizes=[{"w": float(image.width), "h": float(image.height)}],
    )


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _zip_contains(data: bytes, names: tuple[str, ...]) -> bool:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            entries = set(archive.namelist())
            return any(name in entries for name in names)
    except Exception:  # noqa: BLE001
        return False


def _zip_media_count(data: bytes, prefix: str) -> int:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            return sum(1 for name in archive.namelist() if name.startswith(prefix))
    except Exception:  # noqa: BLE001
        return 0
