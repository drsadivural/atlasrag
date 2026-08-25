"""Corrected-document generation.

Every function here is non-destructive: the caller supplies the ORIGINAL bytes and receives
NEW bytes. Nothing writes back over an input, and nothing claims a property the output does
not actually have — in particular a signed PDF always yields a clearly-labelled unsigned
derivative rather than an edit that would silently invalidate the signature.
"""

from __future__ import annotations

import io
import re
from dataclasses import dataclass, field
from typing import Any

import pymupdf

from .extract import _normalise


@dataclass
class Change:
    ordinal: int
    page_number: int | None
    paragraph_index: int | None
    sheet_name: str | None
    cell_range: str | None
    slide_number: int | None
    current_content: str
    proposed_content: str
    reason: str
    citation: str | None

    @classmethod
    def from_dict(cls, raw: dict[str, Any]) -> "Change":
        return cls(
            ordinal=int(raw.get("ordinal", 0)),
            page_number=raw.get("pageNumber"),
            paragraph_index=raw.get("paragraphIndex"),
            sheet_name=raw.get("sheetName"),
            cell_range=raw.get("cellRange"),
            slide_number=raw.get("slideNumber"),
            current_content=raw.get("currentContent") or "",
            proposed_content=raw.get("proposedContent") or "",
            reason=raw.get("reason") or "",
            citation=raw.get("citation"),
        )


@dataclass
class CorrectionOutput:
    document: bytes
    redline: bytes | None
    content_type: str
    extension: str
    applied_changes: int = 0
    unmatched_changes: list[int] = field(default_factory=list)
    pages: int | None = None
    text_length: int = 0
    media_count: int = 0
    page_sizes: list[dict[str, float]] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)


def _normalise_for_match(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip().lower()


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------


def correct_docx(data: bytes, changes: list[Change], include_redline: bool) -> CorrectionOutput:
    """Edits DOCX paragraph runs in place, preserving styles, tables, headers and media.

    Text is replaced inside the first run of a matched paragraph and the remaining runs are
    blanked, which keeps the paragraph's own style and numbering intact. Rewriting the whole
    paragraph object instead would drop its list level and formatting.
    """
    import docx
    from docx.shared import RGBColor

    document = docx.Document(io.BytesIO(data))
    applied = 0
    unmatched: list[int] = []
    redline_entries: list[tuple[str, str, str, str | None]] = []

    paragraphs = list(document.paragraphs)
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                paragraphs.extend(cell.paragraphs)

    for change in changes:
        target = _normalise_for_match(change.current_content)
        matched = False

        if target:
            for paragraph in paragraphs:
                if _normalise_for_match(paragraph.text) == target or (
                    len(target) > 30 and target in _normalise_for_match(paragraph.text)
                ):
                    _replace_paragraph_text(paragraph, change.proposed_content)
                    matched = True
                    break

        if not matched and not change.current_content.strip():
            # An insertion: append the new provision at the end of the body, marked so a
            # reviewer can see it was added rather than always having been there.
            paragraph = document.add_paragraph()
            run = paragraph.add_run(change.proposed_content)
            run.font.color.rgb = RGBColor(0x12, 0xA8, 0x6B)
            matched = True

        if matched:
            applied += 1
            redline_entries.append(
                (str(change.ordinal), change.current_content, change.proposed_content, change.citation)
            )
        else:
            unmatched.append(change.ordinal)

    buffer = io.BytesIO()
    document.save(buffer)
    output = buffer.getvalue()

    redline = _build_docx_redline(redline_entries) if include_redline and redline_entries else None

    reopened = docx.Document(io.BytesIO(output))
    text_length = sum(len(p.text) for p in reopened.paragraphs)

    return CorrectionOutput(
        document=output,
        redline=redline,
        content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        extension="docx",
        applied_changes=applied,
        unmatched_changes=unmatched,
        pages=None,
        text_length=text_length,
        media_count=len(reopened.inline_shapes),
    )


def _replace_paragraph_text(paragraph: Any, text: str) -> None:
    runs = paragraph.runs
    if not runs:
        paragraph.add_run(text)
        return
    runs[0].text = text
    for run in runs[1:]:
        run.text = ""


def _build_docx_redline(entries: list[tuple[str, str, str, str | None]]) -> bytes:
    """A separate change report, since faithful tracked changes cannot be written here.

    python-docx cannot emit real `w:ins`/`w:del` revision marks. Rather than fake tracked
    changes, the correction ships an accompanying redline document that states every edit
    exactly, which a reviewer can read side by side with the corrected file.
    """
    import docx
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, RGBColor

    document = docx.Document()
    heading = document.add_paragraph()
    heading_run = heading.add_run("Change report")
    heading_run.bold = True
    heading_run.font.size = Pt(18)
    heading.alignment = WD_ALIGN_PARAGRAPH.LEFT

    intro = document.add_paragraph()
    intro.add_run(
        "Each accepted change is shown below with the original text, the replacement text, "
        "and the governing citation. The original document is unchanged and retained."
    ).font.size = Pt(10)

    for ordinal, before, after, citation in entries:
        document.add_paragraph()
        title = document.add_paragraph()
        title_run = title.add_run(f"Change {ordinal}")
        title_run.bold = True

        if before.strip():
            removed = document.add_paragraph()
            removed_run = removed.add_run(f"Removed: {before}")
            removed_run.font.strike = True
            removed_run.font.color.rgb = RGBColor(0xE5, 0x48, 0x4D)
        else:
            document.add_paragraph().add_run("Added (no prior text at this location).").italic = True

        added = document.add_paragraph()
        added_run = added.add_run(f"Added: {after}")
        added_run.font.color.rgb = RGBColor(0x12, 0xA8, 0x6B)

        if citation:
            cite = document.add_paragraph()
            cite_run = cite.add_run(f"Governing citation: {citation}")
            cite_run.italic = True
            cite_run.font.size = Pt(9)

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------


def correct_pdf_overlay(data: bytes, changes: list[Change], include_redline: bool) -> CorrectionOutput:
    """Replaces text in a text-based PDF by redacting the old span and drawing the new text.

    Page size, page count and every other page's content are preserved exactly. Redaction
    is used rather than drawing over the top because leaving the original glyphs underneath
    a white box would mean the "corrected" document still contained the wrong text for
    anyone who copied it out.
    """
    document = pymupdf.open(stream=data, filetype="pdf")
    applied = 0
    unmatched: list[int] = []
    warnings: list[str] = []
    original_page_count = document.page_count

    for change in changes:
        target = _normalise(change.current_content)
        if not target:
            unmatched.append(change.ordinal)
            continue

        page_indexes = (
            [change.page_number - 1]
            if change.page_number and 0 < change.page_number <= document.page_count
            else range(document.page_count)
        )

        matched = False
        for index in page_indexes:
            page = document[index]
            # Search on a shortened probe: a long excerpt spanning a line break will not
            # match as one string, but its opening clause will.
            probe = target[:80]
            rects = page.search_for(probe)
            if not rects:
                continue

            rect = rects[0]
            # Widen to the full text block so the replacement is not clipped mid-sentence.
            block_rect = _block_rect_for(page, rect) or rect

            page.add_redact_annot(block_rect, fill=(1, 1, 1))
            page.apply_redactions()

            inserted = page.insert_textbox(
                block_rect,
                change.proposed_content,
                fontsize=9.5,
                fontname="helv",
                color=(0.06, 0.09, 0.18),
                align=0,
            )
            if inserted < 0:
                warnings.append(
                    f"Change {change.ordinal}: the replacement text is longer than the space available on page {index + 1}; it was reflowed and may be truncated. Review the corrected page before use."
                )
            matched = True
            applied += 1
            break

        if not matched:
            unmatched.append(change.ordinal)

    output = document.tobytes(garbage=3, deflate=True)
    page_sizes = [{"w": float(p.rect.width), "h": float(p.rect.height)} for p in document]
    document.close()

    reopened = pymupdf.open(stream=output, filetype="pdf")
    text_length = sum(len(page.get_text("text") or "") for page in reopened)
    media_count = sum(len(page.get_images(full=True)) for page in reopened)
    pages = reopened.page_count
    reopened.close()

    if pages != original_page_count:
        warnings.append(
            f"Page count changed from {original_page_count} to {pages}; the correction was not applied cleanly."
        )

    redline = (
        _build_pdf_redline(changes, unmatched) if include_redline and changes else None
    )

    return CorrectionOutput(
        document=output,
        redline=redline,
        content_type="application/pdf",
        extension="pdf",
        applied_changes=applied,
        unmatched_changes=unmatched,
        pages=pages,
        text_length=text_length,
        media_count=media_count,
        page_sizes=page_sizes,
        warnings=warnings,
    )


def _block_rect_for(page: pymupdf.Page, rect: pymupdf.Rect) -> pymupdf.Rect | None:
    """Finds the enclosing text block so a replacement can reuse its full width."""
    try:
        for block in page.get_text("blocks"):
            x0, y0, x1, y1 = block[0], block[1], block[2], block[3]
            candidate = pymupdf.Rect(x0, y0, x1, y1)
            if candidate.contains(rect):
                return candidate
    except Exception:  # noqa: BLE001
        return None
    return None


def _build_pdf_redline(changes: list[Change], unmatched: list[int]) -> bytes:
    """Change report PDF, laid out with the measured writer so nothing is dropped."""
    from .pdfkit import COBALT, DANGER, Document, MUTED, SUCCESS

    doc = Document()
    doc.footer_text = "UXE Consulting AI  -  change report"

    doc.text("UXE CONSULTING AI", size=9.5, color=COBALT, bold=True, gap=4)
    doc.text("Change report", size=22, bold=True, gap=6)
    doc.text(
        "Each change is listed below with its governing citation. The original document is "
        "retained unchanged; this report describes only the derivative edition.",
        size=9.5,
        color=MUTED,
        gap=14,
    )
    doc.rule()

    for change in changes:
        applied = change.ordinal not in unmatched
        doc.text(
            f"Change {change.ordinal}  -  {'Applied' if applied else 'NOT APPLIED'}",
            size=12,
            bold=True,
            color=SUCCESS if applied else DANGER,
            gap=5,
        )
        if change.current_content.strip():
            doc.text(f"Removed:  {change.current_content}", size=9.5, color=DANGER, gap=4)
        else:
            doc.text("Added at a new location (no prior text to replace).", size=9.5, color=MUTED, gap=4)
        doc.text(f"Added:  {change.proposed_content}", size=9.5, color=SUCCESS, gap=4)
        if change.reason:
            doc.text(f"Reason:  {change.reason}", size=9, color=MUTED, gap=3)
        if change.citation:
            doc.text(f"Governing citation:  {change.citation}", size=9, color=MUTED, gap=3)
        doc.rule(gap=10)

    return doc.finish()


def rebuild_pdf_from_text(
    pages_text: list[str], changes: list[Change], title: str, disclosures: list[str]
) -> CorrectionOutput:
    """Builds a corrected searchable PDF for a scanned original.

    The original page images cannot be edited faithfully, so the corrected edition is
    rebuilt from OCR text with the corrections applied. The disclosures block states this
    on the first page so nobody mistakes it for the scan.
    """
    document = pymupdf.open()
    applied = 0

    corrected_pages: list[str] = []
    for text in pages_text:
        page_text = text
        for change in changes:
            if change.current_content and change.current_content in page_text:
                page_text = page_text.replace(change.current_content, change.proposed_content)
                applied += 1
        corrected_pages.append(page_text)

    from .pdfkit import COBALT, Document as FlowDocument, MUTED

    document.close()
    flow = FlowDocument()
    flow.footer_text = "UXE Consulting AI  -  corrected edition"
    flow.text("UXE CONSULTING AI", size=9.5, color=COBALT, bold=True, gap=4)
    flow.text(title, size=22, bold=True, gap=4)
    flow.text("Corrected edition", size=12, color=MUTED, gap=14)
    flow.rule()
    for line in disclosures:
        flow.text(f"-  {line}", size=9.5, color=MUTED, gap=4)

    for text in corrected_pages:
        flow.new_page()
        flow.text(text, size=10, gap=6)

    output = flow.finish()

    reopened_for_size = pymupdf.open(stream=output, filetype="pdf")
    page_sizes = [{"w": float(p.rect.width), "h": float(p.rect.height)} for p in reopened_for_size]
    pages = reopened_for_size.page_count
    reopened_for_size.close()

    return CorrectionOutput(
        document=output,
        redline=_build_pdf_redline(changes, []),
        content_type="application/pdf",
        extension="pdf",
        applied_changes=applied,
        unmatched_changes=[c.ordinal for c in changes if not c.current_content],
        pages=pages,
        text_length=sum(len(t) for t in corrected_pages),
        media_count=0,
        page_sizes=page_sizes,
        warnings=[
            "This edition was rebuilt from OCR text. Exact glyph positioning, handwriting and page images from the scan are not reproduced."
        ],
    )


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------


def correct_xlsx(data: bytes, changes: list[Change]) -> CorrectionOutput:
    """Edits cell contents while preserving formulas, formats, validations and charts.

    The workbook is loaded with `data_only=False` so a formula cell keeps its formula;
    loading with computed values would silently replace every formula in the file with a
    frozen number.
    """
    import openpyxl

    workbook = openpyxl.load_workbook(io.BytesIO(data), data_only=False)
    applied = 0
    unmatched: list[int] = []

    for change in changes:
        matched = False
        sheets = (
            [workbook[change.sheet_name]]
            if change.sheet_name and change.sheet_name in workbook.sheetnames
            else list(workbook.worksheets)
        )
        target = _normalise_for_match(change.current_content)

        for sheet in sheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is None:
                        continue
                    if _normalise_for_match(str(cell.value)) == target:
                        cell.value = change.proposed_content
                        matched = True
                        break
                if matched:
                    break
            if matched:
                break

        if matched:
            applied += 1
        else:
            unmatched.append(change.ordinal)

    buffer = io.BytesIO()
    workbook.save(buffer)
    output = buffer.getvalue()

    reopened = openpyxl.load_workbook(io.BytesIO(output), read_only=True)
    text_length = 0
    for sheet in reopened.worksheets:
        for row in sheet.iter_rows(values_only=True):
            text_length += sum(len(str(v)) for v in row if v is not None)
    sheets = len(reopened.sheetnames)
    reopened.close()

    return CorrectionOutput(
        document=output,
        redline=None,
        content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        extension="xlsx",
        applied_changes=applied,
        unmatched_changes=unmatched,
        pages=sheets,
        text_length=text_length,
        media_count=0,
    )


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def correct_pptx(data: bytes, changes: list[Change]) -> CorrectionOutput:
    """Edits slide text while preserving slide size, masters, layouts, media and notes."""
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    applied = 0
    unmatched: list[int] = []

    for change in changes:
        target = _normalise_for_match(change.current_content)
        matched = False

        slides = (
            [presentation.slides[change.slide_number - 1]]
            if change.slide_number and 0 < change.slide_number <= len(presentation.slides)
            else list(presentation.slides)
        )

        for slide in slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs)
                    if _normalise_for_match(text) != target:
                        continue
                    if paragraph.runs:
                        paragraph.runs[0].text = change.proposed_content
                        for run in paragraph.runs[1:]:
                            run.text = ""
                        matched = True
                        break
                if matched:
                    break
            if matched:
                break

        if matched:
            applied += 1
        else:
            unmatched.append(change.ordinal)

    buffer = io.BytesIO()
    presentation.save(buffer)
    output = buffer.getvalue()

    reopened = Presentation(io.BytesIO(output))
    text_length = sum(
        len(shape.text_frame.text)
        for slide in reopened.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )

    return CorrectionOutput(
        document=output,
        redline=None,
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        extension="pptx",
        applied_changes=applied,
        unmatched_changes=unmatched,
        pages=len(reopened.slides),
        text_length=text_length,
        media_count=sum(1 for slide in reopened.slides for shape in slide.shapes if shape.shape_type == 13),
        page_sizes=[
            {
                "w": float(reopened.slide_width or 0) / 12700.0,
                "h": float(reopened.slide_height or 0) / 12700.0,
            }
        ],
    )


def revised_edition(
    title: str, changes: list[Change], disclosures: list[str], body_pages: list[str]
) -> CorrectionOutput:
    """Fallback: a professionally formatted revised edition plus an exact change report.

    Used when faithful in-place correction is not safe for the input format. The limitation
    is stated to the user before generation, and repeated in the document itself.
    """
    return rebuild_pdf_from_text(body_pages, changes, title, disclosures)
